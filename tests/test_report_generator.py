"""
Report_Generator 및 CLI 테스트 모듈.

Property 18~21 속성 테스트 및 단위 테스트를 포함한다.
"""

import json
import os
import warnings
from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from hypothesis import given, settings
from hypothesis import strategies as st
from pptx import Presentation

from md_to_pptx.exceptions import FileNotFoundError as CustomFileNotFoundError
from md_to_pptx.main import create_parser, main
from md_to_pptx.models import GenerationResult, SlideContent, SummarizedContent
from md_to_pptx.report_generator import ReportGenerator


# === 헬퍼 함수 ===


def _create_template(tmp_path: Path) -> Path:
    """테스트용 PPTX 템플릿을 생성한다."""
    pptx_path = tmp_path / "template.pptx"
    prs = Presentation()
    prs.save(str(pptx_path))
    return pptx_path


def _create_markdown(tmp_path: Path, filename: str = "test.md") -> Path:
    """테스트용 마크다운 파일을 생성한다."""
    md_path = tmp_path / filename
    md_path.write_text(
        "# 테스트 문서\n\n## 섹션 1\n\n본문 텍스트입니다.\n\n"
        "- 항목 1\n- 항목 2\n",
        encoding="utf-8",
    )
    return md_path


def _mock_bedrock_response(slides: list[dict] | None = None) -> MagicMock:
    """mock된 Bedrock 클라이언트를 생성한다."""
    if slides is None:
        slides = [
            {"title": "표지", "body": ["부제목"], "is_cover": True, "notes": ""},
            {"title": "내용", "body": ["항목 1", "항목 2"], "is_cover": False, "notes": ""},
        ]
    response_json = json.dumps({"slides": slides})
    # Claude 응답 형식
    body_content = json.dumps({
        "content": [{"type": "text", "text": response_json}],
    }).encode("utf-8")

    mock_client = MagicMock()
    mock_response = {"body": MagicMock()}
    mock_response["body"].read.return_value = body_content
    mock_client.invoke_model.return_value = mock_response
    return mock_client


def _generate_report(
    tmp_path: Path,
    md_filename: str = "test.md",
    output_path: str | None = None,
    slides: list[dict] | None = None,
) -> GenerationResult:
    """테스트용 보고서를 생성한다."""
    template_path = _create_template(tmp_path)
    md_path = _create_markdown(tmp_path, md_filename)
    mock_client = _mock_bedrock_response(slides)

    generator = ReportGenerator(bedrock_client=mock_client)
    return generator.generate(
        template_path=str(template_path),
        markdown_path=str(md_path),
        output_path=output_path,
        confirm_overwrite=True,
    )


# === Property 18: PPTX 파일 유효성 라운드트립 ===


class TestProperty18PptxValidity:
    """
    Feature: md-to-pptx-report-generator,
    Property 18: PPTX 파일 유효성 라운드트립

    생성된 PPTX를 python-pptx로 재로드 시 오류 없이 열리고
    슬라이드 수가 일치하는지 검증.
    """

    def test_generated_pptx_can_be_reloaded(self, tmp_path: Path):
        """생성된 PPTX 파일을 python-pptx로 다시 열 수 있다."""
        result = _generate_report(tmp_path)
        # 재로드
        prs = Presentation(result.output_path)
        assert len(prs.slides) == result.slide_count

    def test_slide_count_matches(self, tmp_path: Path):
        """재로드된 PPTX의 슬라이드 수가 GenerationResult와 일치한다."""
        slides_data = [
            {"title": "표지", "body": [], "is_cover": True, "notes": ""},
            {"title": "A", "body": ["a1"], "is_cover": False, "notes": ""},
            {"title": "B", "body": ["b1"], "is_cover": False, "notes": ""},
        ]
        result = _generate_report(tmp_path, slides=slides_data)
        prs = Presentation(result.output_path)
        assert len(prs.slides) == 3
        assert result.slide_count == 3


# === Property 19: 출력 파일명 자동 생성 규칙 ===


class TestProperty19OutputFilename:
    """
    Feature: md-to-pptx-report-generator,
    Property 19: 출력 파일명 자동 생성 규칙

    출력 경로 미지정 시 {마크다운파일명}_report.pptx 형식 검증.
    """

    def test_auto_generated_filename(self, tmp_path: Path):
        """출력 경로 미지정 시 자동 생성된 파일명이 올바르다."""
        result = _generate_report(tmp_path, md_filename="my_document.md")
        expected_name = "my_document_report.pptx"
        assert os.path.basename(result.output_path) == expected_name

    @given(
        name=st.text(
            alphabet=st.characters(
                whitelist_categories=("L", "N"),
            ),
            min_size=1,
            max_size=20,
        ).filter(lambda t: t.strip())
    )
    @settings(max_examples=100)
    def test_auto_filename_format_property(self, name: str):
        """임의의 마크다운 파일명에 대해 출력 파일명 형식이 올바르다."""
        # ReportGenerator._generate_output_path를 직접 테스트
        safe_name = name.strip()
        md_path = f"/tmp/{safe_name}.md"
        output = ReportGenerator._generate_output_path(md_path)
        expected = f"{safe_name}_report.pptx"
        assert os.path.basename(output) == expected


# === Property 20: 생성 결과 완전성 ===


class TestProperty20ResultCompleteness:
    """
    Feature: md-to-pptx-report-generator,
    Property 20: 생성 결과 완전성

    GenerationResult의 output_path, slide_count, elapsed_seconds 검증.
    """

    def test_result_has_valid_output_path(self, tmp_path: Path):
        """GenerationResult의 output_path가 비어있지 않다."""
        result = _generate_report(tmp_path)
        assert result.output_path
        assert len(result.output_path) > 0

    def test_result_has_positive_slide_count(self, tmp_path: Path):
        """GenerationResult의 slide_count가 1 이상이다."""
        result = _generate_report(tmp_path)
        assert result.slide_count >= 1

    def test_result_has_non_negative_elapsed(self, tmp_path: Path):
        """GenerationResult의 elapsed_seconds가 0 이상이다."""
        result = _generate_report(tmp_path)
        assert result.elapsed_seconds >= 0

    def test_output_file_exists(self, tmp_path: Path):
        """생성된 출력 파일이 실제로 존재한다."""
        result = _generate_report(tmp_path)
        assert os.path.isfile(result.output_path)


# === Property 21: CLI 인자 파싱 정확성 ===


class TestProperty21CLIParsing:
    """
    Feature: md-to-pptx-report-generator,
    Property 21: CLI 인자 파싱 정확성

    유효한 CLI 인자 조합에 대해 Namespace 객체 필드가 입력과 일치하는지 검증.
    """

    def test_required_args_parsed(self):
        """필수 인자(template, markdown)가 올바르게 파싱된다."""
        parser = create_parser()
        args = parser.parse_args(["template.pptx", "doc.md"])
        assert args.template == "template.pptx"
        assert args.markdown == "doc.md"
        assert args.output is None
        assert args.font is None

    def test_all_args_parsed(self):
        """모든 인자가 올바르게 파싱된다."""
        parser = create_parser()
        args = parser.parse_args([
            "t.pptx", "d.md", "-o", "out.pptx", "-f", "나눔고딕", "-y"
        ])
        assert args.template == "t.pptx"
        assert args.markdown == "d.md"
        assert args.output == "out.pptx"
        assert args.font == "나눔고딕"
        assert args.yes is True

    @given(
        template=st.text(
            alphabet=st.characters(whitelist_categories=("L", "N")),
            min_size=1,
            max_size=20,
        ).filter(lambda t: t.strip()).map(lambda t: f"{t.strip()}.pptx"),
        markdown=st.text(
            alphabet=st.characters(whitelist_categories=("L", "N")),
            min_size=1,
            max_size=20,
        ).filter(lambda t: t.strip()).map(lambda t: f"{t.strip()}.md"),
    )
    @settings(max_examples=100)
    def test_required_args_property(self, template: str, markdown: str):
        """임의의 필수 인자 조합이 올바르게 파싱된다."""
        parser = create_parser()
        args = parser.parse_args([template, markdown])
        assert args.template == template
        assert args.markdown == markdown


# === 단위 테스트 ===


class TestReportGeneratorUnit:
    """Report_Generator 및 CLI 단위 테스트."""

    def test_cli_is_executable(self):
        """CLI가 실행 가능하다 (요구사항 8.1)."""
        # create_parser가 정상적으로 파서를 반환하는지 확인
        parser = create_parser()
        assert parser is not None
        assert parser.prog == "md-to-pptx"

    def test_missing_required_args_shows_usage(self, capsys):
        """필수 인자 누락 시 사용법이 출력된다 (요구사항 8.3)."""
        parser = create_parser()
        with pytest.raises(SystemExit) as exc_info:
            parser.parse_args([])
        assert exc_info.value.code == 2

    def test_file_not_found_raises_error(self, tmp_path: Path):
        """존재하지 않는 입력 파일에 대해 FileNotFoundError가 발생한다."""
        mock_client = _mock_bedrock_response()
        generator = ReportGenerator(bedrock_client=mock_client)
        with pytest.raises(CustomFileNotFoundError):
            generator.generate(
                template_path=str(tmp_path / "nonexistent.pptx"),
                markdown_path=str(tmp_path / "nonexistent.md"),
            )

    def test_overwrite_raises_file_exists_error(self, tmp_path: Path):
        """출력 파일 중복 시 FileExistsError가 발생한다 (요구사항 6.4)."""
        template_path = _create_template(tmp_path)
        md_path = _create_markdown(tmp_path)
        output_path = tmp_path / "output.pptx"
        output_path.write_text("dummy")  # 기존 파일 생성

        mock_client = _mock_bedrock_response()
        generator = ReportGenerator(bedrock_client=mock_client)

        with pytest.raises(FileExistsError):
            generator.generate(
                template_path=str(template_path),
                markdown_path=str(md_path),
                output_path=str(output_path),
                confirm_overwrite=False,
            )

    def test_cli_main_with_nonexistent_files(self, tmp_path: Path):
        """CLI에서 존재하지 않는 파일 지정 시 오류 코드를 반환한다."""
        exit_code = main([
            str(tmp_path / "no.pptx"),
            str(tmp_path / "no.md"),
        ])
        assert exit_code != 0
