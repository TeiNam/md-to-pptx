"""
Template_Analyzer 테스트 모듈.

속성 기반 테스트(hypothesis)와 단위 테스트(pytest)를 포함한다.
- Property 1: 템플릿 분석 정확성 (요구사항 1.1, 1.2)
- Property 6: 잘못된 입력 파일 오류 처리 (요구사항 1.3, 8.4)
- 단위 테스트: NoPlaceholderWarning 경고 발생 (요구사항 1.4)
"""

import os
import tempfile
import warnings
from pathlib import Path

import pytest
from hypothesis import given, settings
from hypothesis import strategies as st
from pptx import Presentation

from md_to_pptx.exceptions import InvalidFileFormatError, NoPlaceholderWarning
from md_to_pptx.models import PlaceholderType, TemplateInfo
from md_to_pptx.template_analyzer import TemplateAnalyzer, _map_placeholder_type


# === Hypothesis 전략 ===

# python-pptx 기본 Presentation의 레이아웃 인덱스 (0~10, 총 11개)
layout_index_strategy = st.integers(min_value=0, max_value=10)


# === Property 1: 템플릿 분석 정확성 ===


class TestProperty1TemplateAnalysisAccuracy:
    """
    Property 1: 템플릿 분석 정확성

    유효한 PPTX 템플릿에 대해 레이아웃 수와 플레이스홀더 수/유형이 원본과 일치하는지 검증한다.

    Feature: md-to-pptx-report-generator, Property 1: 템플릿 분석 정확성
    Validates: Requirements 1.1, 1.2
    """

    @given(data=st.data())
    @settings(max_examples=100)
    def test_layout_count_matches_original(self, tmp_path_factory, data):
        """분석 결과의 레이아웃 수가 원본 PPTX의 슬라이드 레이아웃 수와 일치해야 한다."""
        # PPTX 템플릿 생성 (기본 Presentation 사용)
        tmp_dir = tmp_path_factory.mktemp("pptx")
        pptx_path = tmp_dir / "template.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        # 원본 레이아웃 수 확인
        expected_layout_count = len(prs.slide_layouts)

        # TemplateAnalyzer로 분석
        analyzer = TemplateAnalyzer()
        result = analyzer.analyze(str(pptx_path))

        # 레이아웃 수 일치 검증
        assert len(result.layouts) == expected_layout_count

    @given(data=st.data())
    @settings(max_examples=100)
    def test_placeholder_count_per_layout_matches(self, tmp_path_factory, data):
        """각 레이아웃의 플레이스홀더 수가 원본과 일치해야 한다."""
        tmp_dir = tmp_path_factory.mktemp("pptx")
        pptx_path = tmp_dir / "template.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        # 원본 레이아웃별 플레이스홀더 수 수집
        expected_counts = []
        for layout in prs.slide_layouts:
            expected_counts.append(len(list(layout.placeholders)))

        # TemplateAnalyzer로 분석
        analyzer = TemplateAnalyzer()
        result = analyzer.analyze(str(pptx_path))

        # 각 레이아웃의 플레이스홀더 수 일치 검증
        for i, layout_info in enumerate(result.layouts):
            assert len(layout_info.placeholders) == expected_counts[i], (
                f"레이아웃 [{i}] '{layout_info.name}': "
                f"기대 {expected_counts[i]}개, 실제 {len(layout_info.placeholders)}개"
            )

    @given(data=st.data())
    @settings(max_examples=100)
    def test_placeholder_types_match_original(self, tmp_path_factory, data):
        """각 플레이스홀더의 유형이 원본과 일치해야 한다."""
        tmp_dir = tmp_path_factory.mktemp("pptx")
        pptx_path = tmp_dir / "template.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        # 임의의 레이아웃 선택
        layout_idx = data.draw(layout_index_strategy)
        original_layout = prs.slide_layouts[layout_idx]

        # 원본 플레이스홀더 유형 수집
        expected_types = []
        for ph in original_layout.placeholders:
            ph_type_value = int(ph.placeholder_format.type)
            expected_types.append(_map_placeholder_type(ph_type_value))

        # TemplateAnalyzer로 분석
        analyzer = TemplateAnalyzer()
        result = analyzer.analyze(str(pptx_path))

        # 해당 레이아웃의 플레이스홀더 유형 일치 검증
        result_layout = result.layouts[layout_idx]
        result_types = [ph.type for ph in result_layout.placeholders]

        assert result_types == expected_types, (
            f"레이아웃 [{layout_idx}] '{result_layout.name}': "
            f"기대 유형 {expected_types}, 실제 유형 {result_types}"
        )

    @given(data=st.data())
    @settings(max_examples=100)
    def test_slide_dimensions_match_original(self, tmp_path_factory, data):
        """슬라이드 너비/높이가 원본과 일치해야 한다."""
        tmp_dir = tmp_path_factory.mktemp("pptx")
        pptx_path = tmp_dir / "template.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        analyzer = TemplateAnalyzer()
        result = analyzer.analyze(str(pptx_path))

        assert result.slide_width == prs.slide_width
        assert result.slide_height == prs.slide_height


# === Property 6: 잘못된 입력 파일 오류 처리 ===


class TestProperty6InvalidInputErrorHandling:
    """
    Property 6: 잘못된 입력 파일 오류 처리

    비-PPTX 파일 또는 존재하지 않는 경로에 대해 적절한 예외 발생을 검증한다.

    Feature: md-to-pptx-report-generator, Property 6: 잘못된 입력 파일 오류 처리
    Validates: Requirements 1.3, 8.4
    """

    @given(
        content=st.binary(min_size=0, max_size=1024),
        extension=st.sampled_from([".txt", ".pdf", ".doc", ".xlsx", ".csv", ".json", ".xml", ".zip"]),
    )
    @settings(max_examples=100)
    def test_non_pptx_extension_raises_error(self, tmp_path_factory, content, extension):
        """비-PPTX 확장자 파일에 대해 InvalidFileFormatError가 발생해야 한다."""
        tmp_dir = tmp_path_factory.mktemp("invalid")
        file_path = tmp_dir / f"test_file{extension}"
        file_path.write_bytes(content)

        analyzer = TemplateAnalyzer()
        with pytest.raises(InvalidFileFormatError):
            analyzer.analyze(str(file_path))

    @given(
        filename=st.text(
            alphabet=st.characters(whitelist_categories=("L", "N"), whitelist_characters="_-"),
            min_size=1,
            max_size=30,
        ),
    )
    @settings(max_examples=100)
    def test_nonexistent_path_raises_error(self, tmp_path_factory, filename):
        """존재하지 않는 PPTX 경로에 대해 InvalidFileFormatError가 발생해야 한다."""
        tmp_dir = tmp_path_factory.mktemp("nonexist")
        nonexistent_path = tmp_dir / f"{filename}.pptx"

        # 파일이 존재하지 않는지 확인
        assert not nonexistent_path.exists()

        analyzer = TemplateAnalyzer()
        with pytest.raises(InvalidFileFormatError):
            analyzer.analyze(str(nonexistent_path))

    @given(content=st.binary(min_size=1, max_size=512))
    @settings(max_examples=100)
    def test_corrupted_pptx_raises_error(self, tmp_path_factory, content):
        """손상된 PPTX 파일(임의 바이트)에 대해 InvalidFileFormatError가 발생해야 한다."""
        tmp_dir = tmp_path_factory.mktemp("corrupted")
        file_path = tmp_dir / "corrupted.pptx"
        file_path.write_bytes(content)

        analyzer = TemplateAnalyzer()
        with pytest.raises(InvalidFileFormatError):
            analyzer.analyze(str(file_path))


# === 단위 테스트: NoPlaceholderWarning 경고 발생 (요구사항 1.4) ===


class TestNoPlaceholderWarning:
    """
    플레이스홀더 없는 템플릿에 대한 NoPlaceholderWarning 경고 발생 확인.

    요구사항 1.4: 템플릿 파일에 플레이스홀더가 하나도 없으면
    "템플릿에 사용 가능한 플레이스홀더가 없습니다" 경고 메시지를 반환한다.
    """

    def test_no_placeholder_template_emits_warning(self, tmp_path):
        """플레이스홀더가 없는 PPTX 템플릿 분석 시 NoPlaceholderWarning이 발생해야 한다."""
        # 플레이스홀더가 없는 빈 PPTX 생성
        # lxml을 사용하여 모든 플레이스홀더를 제거한 PPTX를 만든다
        from pptx.oxml.ns import qn

        prs = Presentation()

        # 모든 슬라이드 레이아웃에서 플레이스홀더 요소를 제거
        for layout in prs.slide_layouts:
            sp_tree = layout.placeholders._element
            # sp (shape) 요소 중 ph (placeholder) 요소를 가진 것을 모두 제거
            shapes_to_remove = []
            for sp in sp_tree.iterchildren(qn("p:sp")):
                ph_elements = sp.findall(f".//{qn('p:ph')}")
                if ph_elements:
                    shapes_to_remove.append(sp)
            for sp in shapes_to_remove:
                sp_tree.remove(sp)

        pptx_path = tmp_path / "no_placeholder.pptx"
        prs.save(str(pptx_path))

        # 경고 캡처
        analyzer = TemplateAnalyzer()
        with warnings.catch_warnings(record=True) as w:
            warnings.simplefilter("always")
            result = analyzer.analyze(str(pptx_path))

            # NoPlaceholderWarning이 발생했는지 확인
            no_ph_warnings = [
                warning for warning in w
                if issubclass(warning.category, NoPlaceholderWarning)
            ]
            assert len(no_ph_warnings) >= 1, (
                "플레이스홀더가 없는 템플릿에서 NoPlaceholderWarning이 발생해야 합니다"
            )

        # 모든 레이아웃의 플레이스홀더가 비어있는지 확인
        for layout_info in result.layouts:
            assert len(layout_info.placeholders) == 0, (
                f"레이아웃 '{layout_info.name}'에 플레이스홀더가 없어야 합니다"
            )
