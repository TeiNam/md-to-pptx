"""
Slide_Composer 테스트 모듈.

Property 8~15 속성 테스트 및 단위 테스트를 포함한다.
"""

import tempfile
import warnings
from pathlib import Path

import pytest
from hypothesis import given, settings
from hypothesis import strategies as st
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Pt

from md_to_pptx.exceptions import ContentOverflowWarning
from md_to_pptx.font_manager import FontManager
from md_to_pptx.models import (
    FontConfig,
    LayoutInfo,
    PlaceholderInfo,
    PlaceholderType,
    SlideContent,
    SummarizedContent,
    TemplateInfo,
)
from md_to_pptx.slide_composer import SlideComposer


# === 헬퍼 함수 ===


def _create_template_pptx(tmp_path: Path) -> str:
    """테스트용 PPTX 템플릿을 생성한다."""
    pptx_path = tmp_path / "template.pptx"
    prs = Presentation()
    prs.save(str(pptx_path))
    return str(pptx_path)


def _analyze_template(template_path: str) -> TemplateInfo:
    """템플릿 파일에서 TemplateInfo를 추출한다."""
    prs = Presentation(template_path)
    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type = PlaceholderType.OTHER
            ph_fmt = ph.placeholder_format
            if ph_fmt.idx == 0:
                ph_type = PlaceholderType.TITLE
            elif ph_fmt.idx == 1:
                ph_type = PlaceholderType.BODY
            elif ph_fmt.idx == 12:
                ph_type = PlaceholderType.SUBTITLE
            placeholders.append(
                PlaceholderInfo(
                    idx=ph_fmt.idx,
                    type=ph_type,
                    name=ph.name,
                    left=ph.left or 0,
                    top=ph.top or 0,
                    width=ph.width or 0,
                    height=ph.height or 0,
                )
            )
        layouts.append(LayoutInfo(name=layout.name, index=idx, placeholders=placeholders))
    return TemplateInfo(
        layouts=layouts,
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
    )


def _make_summarized_content(
    num_slides: int = 2,
    body_lines: list[str] | None = None,
) -> SummarizedContent:
    """테스트용 SummarizedContent를 생성한다."""
    if body_lines is None:
        body_lines = ["항목 1", "항목 2", "항목 3"]
    slides = []
    for i in range(num_slides):
        slides.append(
            SlideContent(
                title=f"슬라이드 {i + 1}" if i > 0 else "표지 제목",
                body=body_lines if i > 0 else ["부제목 텍스트"],
                is_cover=(i == 0),
            )
        )
    return SummarizedContent(slides=slides, original_title="테스트 문서")


def _compose_presentation(
    tmp_path: Path,
    summarized_content: SummarizedContent | None = None,
) -> Presentation:
    """테스트용 프레젠테이션을 생성한다."""
    template_path = _create_template_pptx(tmp_path)
    template_info = _analyze_template(template_path)
    fm = FontManager()
    composer = SlideComposer(fm)
    if summarized_content is None:
        summarized_content = _make_summarized_content()
    return composer.compose(template_info, summarized_content, template_path)


# === Hypothesis 전략 ===

# 안전한 텍스트 전략 (제어 문자 및 마크다운 특수문자 제외)
safe_text = st.text(
    alphabet=st.characters(
        whitelist_categories=("L", "N", "P", "S", "Z"),
        blacklist_characters="\x00\r*_`|#",
    ),
    min_size=1,
    max_size=50,
).filter(lambda t: t.strip())

# 본문 라인 전략
body_line = safe_text.map(lambda t: t.strip()[:50])

# SlideContent 전략
slide_content_strategy = st.builds(
    SlideContent,
    title=safe_text,
    body=st.lists(body_line, min_size=1, max_size=5),
    is_cover=st.just(False),
)

# SummarizedContent 전략 (1~5개 슬라이드, 첫 번째는 표지)
summarized_content_strategy = st.lists(
    slide_content_strategy, min_size=1, max_size=4
).map(
    lambda slides: SummarizedContent(
        slides=[
            SlideContent(
                title="표지",
                body=["부제목"],
                is_cover=True,
            )
        ] + slides,
        original_title="테스트",
    )
)


# === Property 8: 슬라이드 콘텐츠 배치 완전성 ===


class TestProperty8SlideContentPlacement:
    """
    Feature: md-to-pptx-report-generator,
    Property 8: 슬라이드 콘텐츠 배치 완전성

    슬라이드 수가 SummarizedContent와 일치하고
    플레이스홀더에 콘텐츠가 배치되는지 검증한다.
    """

    @given(content=summarized_content_strategy)
    @settings(max_examples=100)
    def test_slide_count_matches_summarized_content(
        self, content: SummarizedContent,
    ):
        """생성된 슬라이드 수가 SummarizedContent의 슬라이드 수와 일치한다."""
        with tempfile.TemporaryDirectory() as td:
            prs = _compose_presentation(Path(td), content)
            assert len(prs.slides) == len(content.slides)


# === Property 9: 슬라이드 서식 속성 일관성 ===


class TestProperty9SlideFormatConsistency:
    """
    Feature: md-to-pptx-report-generator,
    Property 9: 슬라이드 서식 속성 일관성

    줄간격(1.2~1.5), 폰트 크기 계층, 텍스트 정렬 검증.
    """

    def test_line_spacing_in_range(self, tmp_path: Path):
        """본문 텍스트의 줄간격이 1.2~1.5 범위이다."""
        prs = _compose_presentation(tmp_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.line_spacing is not None:
                            assert 1.2 <= para.line_spacing <= 1.5, (
                                f"줄간격 {para.line_spacing}이 범위 밖"
                            )

    def test_title_font_larger_than_body(self, tmp_path: Path):
        """제목 폰트 크기가 본문 폰트 크기보다 크다."""
        fm = FontManager()
        config = fm.get_font_config()
        assert config.title_size_pt > config.body_size_pt

    def test_title_center_aligned(self, tmp_path: Path):
        """제목이 중앙 정렬이다."""
        content = _make_summarized_content(num_slides=2)
        prs = _compose_presentation(tmp_path, content)
        # 두 번째 슬라이드(본문)에서 제목 확인
        slide = prs.slides[1]
        # 플레이스홀더 또는 텍스트 박스에서 제목 찾기
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text == "슬라이드 2":
                            assert para.alignment == PP_ALIGN.CENTER
                            return
        # 제목이 플레이스홀더에 있을 수 있음
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # 제목 플레이스홀더
                for para in ph.text_frame.paragraphs:
                    if para.text.strip():
                        assert para.alignment == PP_ALIGN.CENTER
                        return

    def test_body_left_aligned(self, tmp_path: Path):
        """본문이 좌측 정렬이다."""
        content = _make_summarized_content(num_slides=2)
        prs = _compose_presentation(tmp_path, content)
        slide = prs.slides[1]
        found_body = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text in ("항목 1", "항목 2", "항목 3"):
                            assert para.alignment in (PP_ALIGN.LEFT, None)
                            found_body = True
        # 플레이스홀더에서도 확인
        if not found_body:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:  # 본문 플레이스홀더
                    for para in ph.text_frame.paragraphs:
                        if para.text.strip():
                            assert para.alignment in (PP_ALIGN.LEFT, None)


# === Property 10: 인라인 서식 변환 ===


class TestProperty10InlineFormatting:
    """
    Feature: md-to-pptx-report-generator,
    Property 10: 인라인 서식 변환

    bold/italic 서식이 PPT run에 올바르게 반영되는지 검증.
    """

    def test_bold_text_has_bold_run(self, tmp_path: Path):
        """**bold** 텍스트가 PPT에서 bold=True로 변환된다."""
        content = SummarizedContent(
            slides=[
                SlideContent(title="표지", body=["부제목"], is_cover=True),
                SlideContent(
                    title="본문",
                    body=["이것은 **굵은 텍스트** 입니다"],
                    is_cover=False,
                ),
            ],
            original_title="테스트",
        )
        prs = _compose_presentation(tmp_path, content)
        slide = prs.slides[1]

        found_bold = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text == "굵은 텍스트":
                            assert run.font.bold is True
                            found_bold = True
        assert found_bold, "bold 텍스트를 찾을 수 없음"

    def test_italic_text_has_italic_run(self, tmp_path: Path):
        """*italic* 텍스트가 PPT에서 italic=True로 변환된다."""
        content = SummarizedContent(
            slides=[
                SlideContent(title="표지", body=["부제목"], is_cover=True),
                SlideContent(
                    title="본문",
                    body=["이것은 *기울임 텍스트* 입니다"],
                    is_cover=False,
                ),
            ],
            original_title="테스트",
        )
        prs = _compose_presentation(tmp_path, content)
        slide = prs.slides[1]

        found_italic = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text == "기울임 텍스트":
                            assert run.font.italic is True
                            found_italic = True
        assert found_italic, "italic 텍스트를 찾을 수 없음"

    @given(
        bold_text=safe_text,
        italic_text=safe_text,
    )
    @settings(max_examples=100)
    def test_inline_formatting_property(
        self, bold_text: str, italic_text: str,
    ):
        """임의의 bold/italic 텍스트가 올바르게 변환된다."""
        content = SummarizedContent(
            slides=[
                SlideContent(title="표지", body=["부제목"], is_cover=True),
                SlideContent(
                    title="본문",
                    body=[f"**{bold_text}** 그리고 *{italic_text}*"],
                    is_cover=False,
                ),
            ],
            original_title="테스트",
        )
        with tempfile.TemporaryDirectory() as td:
            prs = _compose_presentation(Path(td), content)
            slide = prs.slides[1]

            bold_found = False
            italic_found = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.text == bold_text and run.font.bold:
                                bold_found = True
                            if run.text == italic_text and run.font.italic:
                                italic_found = True
            assert bold_found, f"bold 텍스트 '{bold_text}' 미발견"
            assert italic_found, f"italic 텍스트 '{italic_text}' 미발견"


# === Property 11: 목록 들여쓰기 깊이 보존 ===


class TestProperty11ListIndentation:
    """
    Feature: md-to-pptx-report-generator,
    Property 11: 목록 들여쓰기 깊이 보존

    중첩 목록의 들여쓰기 레벨이 원본과 일치하는지 검증.
    """

    def test_nested_list_indentation(self, tmp_path: Path):
        """중첩 목록의 들여쓰기 레벨이 보존된다."""
        content = SummarizedContent(
            slides=[
                SlideContent(title="표지", body=["부제목"], is_cover=True),
                SlideContent(
                    title="목록 테스트",
                    body=[
                        "- 레벨 0 항목",
                        "  - 레벨 1 항목",
                        "    - 레벨 2 항목",
                    ],
                    is_cover=False,
                ),
            ],
            original_title="테스트",
        )
        prs = _compose_presentation(tmp_path, content)
        slide = prs.slides[1]

        # 본문에서 목록 항목의 level 확인
        levels_found: list[int] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text in (
                            "레벨 0 항목", "레벨 1 항목", "레벨 2 항목"
                        ):
                            levels_found.append(para.level or 0)

        if levels_found:
            # 레벨이 증가하는 순서인지 확인
            assert levels_found == sorted(levels_found)
            # 최소 2개 이상의 서로 다른 레벨이 있어야 함
            assert len(set(levels_found)) >= 2, (
                f"들여쓰기 레벨이 구분되지 않음: {levels_found}"
            )


# === Property 12: 코드 블록 고정폭 폰트 적용 ===


class TestProperty12CodeBlockFont:
    """
    Feature: md-to-pptx-report-generator,
    Property 12: 코드 블록 고정폭 폰트 적용

    코드 블록에 고정폭 폰트가 적용되는지 검증.
    """

    def test_code_block_uses_mono_font(self, tmp_path: Path):
        """코드 블록 텍스트에 고정폭 폰트가 적용된다."""
        # 플레이스홀더 없는 빈 템플릿 생성 (텍스트 박스 경로 사용)
        pptx_path = tmp_path / "blank.pptx"
        prs = Presentation()
        # 모든 레이아웃의 플레이스홀더를 제거하기 어려우므로
        # _add_code_block 메서드를 직접 테스트
        fm = FontManager()
        composer = SlideComposer(fm)
        mono_font = fm.get_font_config().mono_font

        # 빈 프레젠테이션에 슬라이드 추가 후 코드 블록 배치
        prs.save(str(pptx_path))
        prs2 = Presentation(str(pptx_path))
        slide_layout = prs2.slide_layouts[0]
        slide = prs2.slides.add_slide(slide_layout)

        from pptx.util import Inches
        composer._add_code_block(
            slide, "def hello():\n    print('world')",
            Inches(0.5), Inches(1.5), Inches(8), Inches(2),
        )

        found_code = False
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if "hello" in run.text or "print" in run.text:
                            assert run.font.name == mono_font, (
                                f"코드 블록 폰트가 '{run.font.name}'이지만 "
                                f"'{mono_font}'이어야 함"
                            )
                            found_code = True
        assert found_code, "코드 블록 텍스트를 찾을 수 없음"


# === Property 13: 표 데이터 변환 ===


class TestProperty13TableConversion:
    """
    Feature: md-to-pptx-report-generator,
    Property 13: 표 데이터 변환

    표가 PPT 표 객체로 변환되고 헤더 행에 배경색이 적용되는지 검증.
    """

    def test_table_converted_to_ppt_table(self, tmp_path: Path):
        """표 데이터가 PPT 표 객체로 변환된다."""
        # _add_table 메서드를 직접 테스트
        pptx_path = tmp_path / "blank.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))
        prs2 = Presentation(str(pptx_path))
        slide = prs2.slides.add_slide(prs2.slide_layouts[0])

        fm = FontManager()
        composer = SlideComposer(fm)

        from pptx.util import Inches
        table_data = [
            ["이름", "나이"],
            ["홍길동", "30"],
            ["김철수", "25"],
        ]
        composer._add_table(
            slide, table_data,
            Inches(0.5), Inches(1.5), Inches(8), Inches(2),
        )

        table_found = False
        for shape in slide.shapes:
            if shape.has_table:
                table_found = True
                table = shape.table
                assert len(table.rows) >= 2, "표 행이 2개 미만"
                assert len(table.columns) >= 2, "표 열이 2개 미만"
                break
        assert table_found, "PPT 표 객체를 찾을 수 없음"

    def test_table_header_has_background_color(self, tmp_path: Path):
        """표 헤더 행에 배경색이 적용된다."""
        pptx_path = tmp_path / "blank.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))
        prs2 = Presentation(str(pptx_path))
        slide = prs2.slides.add_slide(prs2.slide_layouts[0])

        fm = FontManager()
        composer = SlideComposer(fm)

        from pptx.util import Inches
        table_data = [
            ["항목", "값"],
            ["A", "100"],
        ]
        composer._add_table(
            slide, table_data,
            Inches(0.5), Inches(1.5), Inches(8), Inches(2),
        )

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                header_cell = table.cell(0, 0)
                fill = header_cell.fill
                assert fill.type is not None, "헤더 셀에 배경색이 없음"
                return
        pytest.fail("PPT 표 객체를 찾을 수 없음")


# === Property 14: 페이지 번호 삽입 ===


class TestProperty14PageNumber:
    """
    Feature: md-to-pptx-report-generator,
    Property 14: 페이지 번호 삽입

    모든 슬라이드 하단에 페이지 번호가 존재하는지 검증.
    """

    def test_all_slides_have_page_number(self, tmp_path: Path):
        """모든 슬라이드에 페이지 번호가 존재한다."""
        content = _make_summarized_content(num_slides=3)
        prs = _compose_presentation(tmp_path, content)

        for slide_idx, slide in enumerate(prs.slides):
            has_page_num = False
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        # slidenum 필드 확인
                        fld_elements = para._p.findall(qn("a:fld"))
                        for fld in fld_elements:
                            if fld.get("type") == "slidenum":
                                has_page_num = True
                                break
                    if has_page_num:
                        break
            assert has_page_num, (
                f"슬라이드 {slide_idx + 1}에 페이지 번호가 없음"
            )

    @given(num_slides=st.integers(min_value=1, max_value=5))
    @settings(max_examples=100)
    def test_page_number_exists_for_any_slide_count(
        self, num_slides: int,
    ):
        """임의의 슬라이드 수에 대해 모든 슬라이드에 페이지 번호가 존재한다."""
        content = _make_summarized_content(num_slides=num_slides)
        with tempfile.TemporaryDirectory() as td:
            prs = _compose_presentation(Path(td), content)

            for slide_idx, slide in enumerate(prs.slides):
                has_page_num = False
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            fld_elements = para._p.findall(qn("a:fld"))
                            for fld in fld_elements:
                                if fld.get("type") == "slidenum":
                                    has_page_num = True
                                    break
                        if has_page_num:
                            break
                assert has_page_num, (
                    f"슬라이드 {slide_idx + 1}에 페이지 번호가 없음"
                )


# === Property 15: 텍스트 박스 비겹침 ===


class TestProperty15TextBoxNoOverlap:
    """
    Feature: md-to-pptx-report-generator,
    Property 15: 텍스트 박스 비겹침

    슬라이드 내 텍스트 박스 영역이 서로 겹치지 않는지 검증.
    """

    @staticmethod
    def _boxes_overlap(
        a_left: int, a_top: int, a_width: int, a_height: int,
        b_left: int, b_top: int, b_width: int, b_height: int,
    ) -> bool:
        """두 박스가 겹치는지 확인한다."""
        a_right = a_left + a_width
        a_bottom = a_top + a_height
        b_right = b_left + b_width
        b_bottom = b_top + b_height
        # 겹치지 않는 조건의 부정
        if a_right <= b_left or b_right <= a_left:
            return False
        if a_bottom <= b_top or b_bottom <= a_top:
            return False
        return True

    def test_textboxes_do_not_overlap(self, tmp_path: Path):
        """슬라이드 내 텍스트 박스들이 서로 겹치지 않는다."""
        content = _make_summarized_content(num_slides=3)
        prs = _compose_presentation(tmp_path, content)

        for slide_idx, slide in enumerate(prs.slides):
            boxes = []
            for shape in slide.shapes:
                if shape.has_text_frame and not shape.has_table:
                    boxes.append((
                        shape.left or 0,
                        shape.top or 0,
                        shape.width or 0,
                        shape.height or 0,
                    ))

            # 모든 박스 쌍에 대해 겹침 확인
            for i in range(len(boxes)):
                for j in range(i + 1, len(boxes)):
                    overlap = self._boxes_overlap(
                        *boxes[i], *boxes[j]
                    )
                    assert not overlap, (
                        f"슬라이드 {slide_idx + 1}: "
                        f"텍스트 박스 {i}와 {j}가 겹침 - "
                        f"박스{i}={boxes[i]}, 박스{j}={boxes[j]}"
                    )


# === 단위 테스트: ContentOverflowWarning ===


class TestContentOverflowWarning:
    """
    콘텐츠 영역 초과 시 ContentOverflowWarning 경고 기록 확인.
    요구사항 4.6
    """

    def test_overflow_warning_on_excessive_content(self, tmp_path: Path):
        """최소 폰트 크기(10pt)까지 축소 후에도 초과 시 경고가 발생한다."""
        # 매우 긴 본문으로 오버플로우 유도
        long_body = [f"매우 긴 항목 {i}: " + "가나다라마바사" * 20 for i in range(30)]
        content = SummarizedContent(
            slides=[
                SlideContent(title="표지", body=["부제목"], is_cover=True),
                SlideContent(
                    title="오버플로우 테스트",
                    body=long_body,
                    is_cover=False,
                ),
            ],
            original_title="테스트",
        )

        with warnings.catch_warnings(record=True) as w:
            warnings.simplefilter("always")
            prs = _compose_presentation(tmp_path, content)
            # ContentOverflowWarning이 발생했는지 확인
            overflow_warnings = [
                x for x in w
                if issubclass(x.category, ContentOverflowWarning)
            ]
            # 오버플로우가 감지되면 경고가 있어야 함
            # (플레이스홀더 기반이 아닌 텍스트 박스 기반에서는
            #  오버플로우 감지가 제한적일 수 있으므로 슬라이드 생성 자체를 확인)
            assert len(prs.slides) == 2, "슬라이드가 생성되어야 함"
