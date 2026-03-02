"""
PPT 템플릿 분석 모듈.

PPTX 템플릿 파일의 슬라이드 레이아웃, 플레이스홀더 유형/위치를 분석하여
구조화된 TemplateInfo 데이터로 반환한다.
"""

import logging
import os
import warnings

from pptx import Presentation
from pptx.exc import PackageNotFoundError

from md_to_pptx.exceptions import InvalidFileFormatError, NoPlaceholderWarning
from md_to_pptx.models import (
    LayoutInfo,
    PlaceholderInfo,
    PlaceholderType,
    TemplateInfo,
)

logger = logging.getLogger(__name__)

# python-pptx placeholder type 값 → PlaceholderType 매핑
# PP_PLACEHOLDER_TYPE enum의 정수 값을 기준으로 매핑한다
_PLACEHOLDER_TYPE_MAP: dict[int, PlaceholderType] = {
    1: PlaceholderType.TITLE,         # TITLE
    2: PlaceholderType.BODY,          # BODY
    3: PlaceholderType.TITLE,         # CENTER_TITLE → TITLE로 매핑
    4: PlaceholderType.SUBTITLE,      # SUBTITLE
    7: PlaceholderType.BODY,          # OBJECT → BODY로 매핑
    18: PlaceholderType.IMAGE,        # PICTURE → IMAGE로 매핑
}


def _map_placeholder_type(pptx_type_value: int) -> PlaceholderType:
    """python-pptx 플레이스홀더 타입 값을 PlaceholderType enum으로 변환한다."""
    return _PLACEHOLDER_TYPE_MAP.get(pptx_type_value, PlaceholderType.OTHER)


class TemplateAnalyzer:
    """PPTX 템플릿 파일을 분석하여 레이아웃 구조를 추출하는 클래스."""

    def analyze(self, template_path: str) -> TemplateInfo:
        """
        PPTX 템플릿 파일을 분석하여 레이아웃 정보를 반환한다.

        Args:
            template_path: PPTX 템플릿 파일 경로

        Returns:
            TemplateInfo: 레이아웃 목록과 플레이스홀더 정보

        Raises:
            InvalidFileFormatError: 유효하지 않은 파일 형식이거나 파일이 존재하지 않는 경우
        """
        logger.info("템플릿 분석 시작: %s", template_path)

        # 파일 확장자 검증
        if not template_path.lower().endswith(".pptx"):
            logger.error("지원하지 않는 파일 형식: %s", template_path)
            raise InvalidFileFormatError(template_path)

        # 파일 존재 여부 검증
        if not os.path.exists(template_path):
            logger.error("파일을 찾을 수 없음: %s", template_path)
            raise InvalidFileFormatError(template_path)

        # PPTX 파일 로드
        try:
            prs = Presentation(template_path)
        except PackageNotFoundError:
            logger.error("PPTX 파일 로드 실패: %s", template_path)
            raise InvalidFileFormatError(template_path)
        except Exception:
            logger.error("PPTX 파일 로드 중 예상치 못한 오류: %s", template_path)
            raise InvalidFileFormatError(template_path)

        logger.info("PPTX 파일 로드 완료, 레이아웃 분석 중...")

        # 슬라이드 레이아웃 분석
        layouts: list[LayoutInfo] = []
        total_placeholder_count = 0

        for index, slide_layout in enumerate(prs.slide_layouts):
            placeholders = self._extract_placeholders(slide_layout)
            total_placeholder_count += len(placeholders)

            layout_info = LayoutInfo(
                name=slide_layout.name,
                index=index,
                placeholders=placeholders,
            )
            layouts.append(layout_info)

            logger.debug(
                "레이아웃 [%d] '%s': 플레이스홀더 %d개",
                index,
                slide_layout.name,
                len(placeholders),
            )

        # 플레이스홀더가 하나도 없는 경우 경고 발행
        if total_placeholder_count == 0:
            logger.warning("템플릿에 사용 가능한 플레이스홀더가 없습니다")
            warnings.warn(
                NoPlaceholderWarning(),
                stacklevel=2,
            )

        template_info = TemplateInfo(
            layouts=layouts,
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
        )

        logger.info(
            "템플릿 분석 완료: 레이아웃 %d개, 총 플레이스홀더 %d개",
            len(layouts),
            total_placeholder_count,
        )

        return template_info

    def _extract_placeholders(self, slide_layout) -> list[PlaceholderInfo]:
        """슬라이드 레이아웃에서 플레이스홀더 정보를 추출한다."""
        placeholders: list[PlaceholderInfo] = []

        for ph in slide_layout.placeholders:
            ph_fmt = ph.placeholder_format
            ph_type_value = int(ph_fmt.type)
            mapped_type = _map_placeholder_type(ph_type_value)

            placeholder_info = PlaceholderInfo(
                idx=ph_fmt.idx,
                type=mapped_type,
                name=ph.name,
                left=ph.left,
                top=ph.top,
                width=ph.width,
                height=ph.height,
            )
            placeholders.append(placeholder_info)

        return placeholders
