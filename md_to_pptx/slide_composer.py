"""
슬라이드 구성 모듈.

요약된 콘텐츠를 PPT 템플릿 레이아웃에 맞게 슬라이드에 배치하고
서식(줄간격, 불릿포인트, 폰트 크기, 정렬 등)을 적용한다.
"""

import logging
import re
import uuid
import warnings

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from md_to_pptx.exceptions import ContentOverflowWarning
from md_to_pptx.font_manager import FontManager
from md_to_pptx.models import (
    FontConfig,
    LayoutInfo,
    PlaceholderInfo,
    PlaceholderType,
    SlideContent,
    SummarizedContent,
    TemplateInfo,
)

logger = logging.getLogger(__name__)

# 상수 정의
_LINE_SPACING = 1.3  # 본문 줄간격 (1.2~1.5 범위)
_FONT_SHRINK_STEP_PT = 2  # 폰트 축소 단계 (pt)
_MIN_TEXTBOX_GAP_EMU = Inches(0.1)  # 텍스트 박스 간 최소 간격
_CODE_BG_COLOR = RGBColor(0xF2, 0xF2, 0xF2)  # 코드 블록 배경색
_TABLE_HEADER_BG = RGBColor(0x4F, 0x81, 0xBD)  # 표 헤더 배경색
_TABLE_HEADER_FONT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)  # 표 헤더 폰트 색상
_PAGE_NUM_HEIGHT = Inches(0.3)  # 페이지 번호 영역 높이
_PAGE_NUM_WIDTH = Inches(1.5)  # 페이지 번호 영역 너비
_MARGIN = Inches(0.5)  # 기본 여백
_BULLET_CHARS = ["•", "–", "◦", "▪"]  # 깊이별 불릿 문자
_INDENT_EMU = Emu(228600)  # 들여쓰기 단위 (0.25인치)



class SlideComposer:
    """요약 콘텐츠를 템플릿에 배치하여 프레젠테이션을 생성하는 클래스."""

    def __init__(self, font_manager: FontManager) -> None:
        """
        SlideComposer를 초기화한다.

        Args:
            font_manager: 폰트 설정을 관리하는 FontManager 인스턴스
        """
        self._font_manager = font_manager
        self._font_config: FontConfig = font_manager.get_font_config()
        logger.info("SlideComposer 초기화 완료")

    def compose(
        self,
        template_info: TemplateInfo,
        summarized_content: SummarizedContent,
        template_path: str,
    ) -> Presentation:
        """
        요약 콘텐츠를 템플릿에 배치하여 Presentation 객체를 생성한다.

        Args:
            template_info: 템플릿 분석 결과
            summarized_content: 요약된 슬라이드 콘텐츠
            template_path: 원본 템플릿 파일 경로

        Returns:
            python-pptx Presentation 객체
        """
        logger.info(
            "슬라이드 구성 시작: %d개 슬라이드",
            len(summarized_content.slides),
        )
        prs = Presentation(template_path)

        # 기존 슬라이드 제거 (템플릿에 포함된 샘플 슬라이드)
        while len(prs.slides) > 0:
            rId = prs.slides._sldIdLst[0].get(qn("r:id"))
            prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

        for idx, slide_content in enumerate(summarized_content.slides):
            logger.info(
                "슬라이드 %d/%d 생성 중: '%s'",
                idx + 1,
                len(summarized_content.slides),
                slide_content.title,
            )
            layout = self._select_layout(template_info, slide_content)
            slide = prs.slides.add_slide(prs.slide_layouts[layout.index])
            self._populate_slide(
                slide, slide_content, layout, template_info, idx
            )
            self._add_page_number(slide, template_info, idx + 1)

        logger.info("슬라이드 구성 완료: 총 %d개", len(prs.slides))
        return prs

    def _select_layout(
        self, template_info: TemplateInfo, slide_content: SlideContent
    ) -> LayoutInfo:
        """슬라이드 콘텐츠에 적합한 레이아웃을 선택한다."""
        if not template_info.layouts:
            # 레이아웃이 없으면 인덱스 0 사용
            return LayoutInfo(name="Default", index=0)

        if slide_content.is_cover:
            # 표지: 제목 슬라이드 레이아웃 (보통 인덱스 0)
            for layout in template_info.layouts:
                title_ph = [
                    p for p in layout.placeholders
                    if p.type == PlaceholderType.TITLE
                ]
                subtitle_ph = [
                    p for p in layout.placeholders
                    if p.type == PlaceholderType.SUBTITLE
                ]
                if title_ph and subtitle_ph:
                    return layout
            return template_info.layouts[0]

        # 본문: 제목+본문 레이아웃 선택
        for layout in template_info.layouts:
            title_ph = [
                p for p in layout.placeholders
                if p.type == PlaceholderType.TITLE
            ]
            body_ph = [
                p for p in layout.placeholders
                if p.type == PlaceholderType.BODY
            ]
            if title_ph and body_ph:
                return layout

        # 적합한 레이아웃이 없으면 첫 번째 사용
        return template_info.layouts[0]

    def _populate_slide(
        self,
        slide,
        slide_content: SlideContent,
        layout: LayoutInfo,
        template_info: TemplateInfo,
        slide_index: int,
    ) -> None:
        """슬라이드에 콘텐츠를 배치한다."""
        # 플레이스홀더에 제목 배치
        title_placed = False
        body_placed = False

        for ph in slide.placeholders:
            ph_idx = ph.placeholder_format.idx
            ph_info = self._find_placeholder_info(layout, ph_idx)

            if ph_info and ph_info.type == PlaceholderType.TITLE:
                self._set_title_text(ph, slide_content.title)
                title_placed = True
            elif ph_info and ph_info.type == PlaceholderType.SUBTITLE:
                if slide_content.is_cover and slide_content.body:
                    self._set_body_text(ph, slide_content.body, is_subtitle=True)
                    body_placed = True
            elif ph_info and ph_info.type == PlaceholderType.BODY:
                self._set_body_text(ph, slide_content.body)
                body_placed = True

        # 플레이스홀더가 없으면 텍스트 박스로 직접 배치
        sw = template_info.slide_width
        sh = template_info.slide_height

        if not title_placed:
            self._add_title_textbox(slide, slide_content.title, sw, sh)

        if not body_placed and slide_content.body:
            self._add_body_textbox(
                slide, slide_content, sw, sh, slide_index
            )

        # 발표자 노트 추가
        if slide_content.notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_content.notes

    def _find_placeholder_info(
        self, layout: LayoutInfo, ph_idx: int
    ) -> PlaceholderInfo | None:
        """레이아웃에서 플레이스홀더 정보를 찾는다."""
        for ph_info in layout.placeholders:
            if ph_info.idx == ph_idx:
                return ph_info
        return None

    def _set_title_text(self, placeholder, title: str) -> None:
        """플레이스홀더에 제목 텍스트를 설정한다."""
        tf = placeholder.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._add_formatted_runs(p, title, is_title=True)

    def _set_body_text(
        self, placeholder, body_lines: list[str], is_subtitle: bool = False
    ) -> None:
        """플레이스홀더에 본문 텍스트를 설정한다."""
        tf = placeholder.text_frame
        tf.clear()
        tf.word_wrap = True

        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = _LINE_SPACING

            # 목록 항목 처리
            depth, clean_line = self._parse_list_item(line)
            if depth >= 0:
                p.level = depth
                self._set_bullet(p, depth)
                self._add_formatted_runs(p, clean_line, is_title=False)
            elif self._is_code_block(line):
                self._format_code_paragraph(p, line)
            elif self._is_table_line(line):
                # 표는 별도 처리 (플레이스홀더 내에서는 텍스트로 표시)
                self._add_formatted_runs(p, line, is_title=False)
            else:
                self._add_formatted_runs(
                    p, line, is_title=is_subtitle
                )

    def _add_title_textbox(
        self, slide, title: str, slide_width: int, slide_height: int
    ) -> None:
        """제목 텍스트 박스를 슬라이드에 추가한다."""
        left = _MARGIN
        top = _MARGIN
        width = slide_width - 2 * _MARGIN
        height = Inches(1.0)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        self._add_formatted_runs(p, title, is_title=True)

    def _add_body_textbox(
        self,
        slide,
        slide_content: SlideContent,
        slide_width: int,
        slide_height: int,
        slide_index: int,
    ) -> None:
        """본문 텍스트 박스를 슬라이드에 추가한다."""
        left = _MARGIN
        # 제목 아래 + 간격
        top = _MARGIN + Inches(1.0) + _MIN_TEXTBOX_GAP_EMU
        width = slide_width - 2 * _MARGIN
        # 페이지 번호 영역 확보
        height = (
            slide_height - top - _PAGE_NUM_HEIGHT - _MARGIN
        )

        # 표 데이터와 코드 블록을 분리하여 처리
        regular_lines: list[str] = []
        table_data: list[list[list[str]]] = []
        code_blocks: list[str] = []

        i = 0
        body = slide_content.body
        while i < len(body):
            line = body[i]
            if self._is_table_line(line):
                table = self._extract_table_data(body, i)
                if table:
                    table_data.append(table[0])
                    i = table[1]
                    continue
            elif self._is_code_block_marker(line):
                code, end_idx = self._extract_code_block(body, i)
                if code is not None:
                    code_blocks.append(code)
                    i = end_idx
                    continue
            regular_lines.append(line)
            i += 1

        # 본문 텍스트 배치
        current_top = top
        if regular_lines:
            # 표/코드가 있으면 텍스트와 표/코드 간 공간을 비율로 분배
            if table_data or code_blocks:
                # 텍스트 라인 수 기반 비율 계산
                text_lines = len(regular_lines)
                table_rows = sum(len(t) for t in table_data)
                total_items = text_lines + table_rows + len(code_blocks) * 3
                text_ratio = max(0.2, min(0.6, text_lines / total_items)) if total_items > 0 else 0.5
                body_height = int(height * text_ratio)
            else:
                body_height = height

            txBox = slide.shapes.add_textbox(
                left, current_top, width, body_height
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            font_size_pt = self._font_config.body_size_pt
            overflow = self._fill_body_text(
                tf, regular_lines, font_size_pt
            )

            # 오버플로우 시 폰트 축소
            if overflow:
                font_size_pt = self._shrink_font_to_fit(
                    tf, regular_lines, font_size_pt, slide_index
                )

            current_top += body_height + _MIN_TEXTBOX_GAP_EMU

        # 코드 블록 배치
        for code in code_blocks:
            if current_top >= slide_height - _PAGE_NUM_HEIGHT - _MARGIN:
                break
            remaining = (
                slide_height - current_top - _PAGE_NUM_HEIGHT - _MARGIN
            )
            code_height = min(Inches(1.5), remaining)
            self._add_code_block(
                slide, code, left, current_top, width, code_height
            )
            current_top += code_height + _MIN_TEXTBOX_GAP_EMU

        # 표 배치
        for table in table_data:
            if current_top >= slide_height - _PAGE_NUM_HEIGHT - _MARGIN:
                break
            remaining = (
                slide_height - current_top - _PAGE_NUM_HEIGHT - _MARGIN
            )
            # 행 수에 비례한 표 높이 계산 (행당 약 0.35인치, 최소 1인치)
            row_count = len(table)
            estimated_height = max(Inches(1.0), Inches(0.35) * row_count)
            table_height = min(estimated_height, remaining)
            self._add_table(
                slide, table, left, current_top, width, table_height
            )
            current_top += table_height + _MIN_TEXTBOX_GAP_EMU

    def _fill_body_text(
        self, tf, lines: list[str], font_size_pt: int
    ) -> bool:
        """텍스트 프레임에 본문 텍스트를 채운다. 오버플로우 여부를 반환한다."""
        tf.clear()
        total_height_est = 0
        line_height_est = Pt(font_size_pt) * _LINE_SPACING * 1.2

        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = _LINE_SPACING

            depth, clean_line = self._parse_list_item(line)
            if depth >= 0:
                p.level = depth
                self._set_bullet(p, depth)
                self._add_formatted_runs(p, clean_line, is_title=False,
                                         font_size_pt=font_size_pt)
            else:
                self._add_formatted_runs(p, line, is_title=False,
                                         font_size_pt=font_size_pt)

            total_height_est += line_height_est

        # 텍스트 프레임 높이와 비교하여 오버플로우 추정
        # 1) shape의 height 속성 사용 (텍스트 박스)
        frame_height = 0
        try:
            shape = tf._txBody.getparent()
            if shape is not None:
                # spPr/a:xfrm/a:ext 에서 cy 추출
                xfrm = shape.find(qn('p:spPr'))
                if xfrm is None:
                    xfrm = shape.find(qn('a:xfrm'))
                if xfrm is not None:
                    ext = xfrm.find(qn('a:ext'))
                    if ext is not None:
                        frame_height = int(ext.get('cy', '0'))
                # spPr 내부의 xfrm 탐색
                if frame_height == 0:
                    sp_pr = shape.find(qn('p:spPr'))
                    if sp_pr is not None:
                        a_xfrm = sp_pr.find(qn('a:xfrm'))
                        if a_xfrm is not None:
                            ext = a_xfrm.find(qn('a:ext'))
                            if ext is not None:
                                frame_height = int(ext.get('cy', '0'))
        except (AttributeError, TypeError):
            pass

        if frame_height > 0 and total_height_est > frame_height:
            return True
        return False

    def _shrink_font_to_fit(
        self,
        tf,
        lines: list[str],
        current_size_pt: int,
        slide_index: int,
    ) -> int:
        """폰트 크기를 단계적으로 축소하여 영역에 맞춘다."""
        min_size = self._font_config.min_size_pt

        while current_size_pt > min_size:
            current_size_pt = max(
                current_size_pt - _FONT_SHRINK_STEP_PT, min_size
            )
            logger.debug(
                "슬라이드 %d: 폰트 크기 %dpt로 축소",
                slide_index + 1,
                current_size_pt,
            )
            overflow = self._fill_body_text(tf, lines, current_size_pt)
            if not overflow:
                return current_size_pt

        # 최소 크기에서도 오버플로우 → 경고 기록
        logger.warning(
            "슬라이드 %d: 최소 폰트 크기(%dpt)에서도 콘텐츠 초과",
            slide_index + 1,
            min_size,
        )
        warnings.warn(ContentOverflowWarning(slide_index + 1))
        return min_size

    def _add_formatted_runs(
        self,
        paragraph,
        text: str,
        is_title: bool = False,
        font_size_pt: int | None = None,
    ) -> None:
        """인라인 서식(bold/italic)을 파싱하여 Run으로 추가한다."""
        if font_size_pt is None:
            font_size_pt = (
                self._font_config.title_size_pt
                if is_title
                else self._font_config.body_size_pt
            )

        # 인라인 서식 파싱: **bold**, *italic*, ***bold+italic***, `code`
        segments = self._parse_inline_formatting(text)

        for seg_text, bold, italic, is_code in segments:
            run = paragraph.add_run()
            run.text = seg_text
            run.font.size = Pt(font_size_pt)

            if is_code:
                run.font.name = self._font_config.mono_font
            else:
                run.font.name = self._font_config.korean_font

            if bold:
                run.font.bold = True
            if italic:
                run.font.italic = True

    @staticmethod
    def _parse_inline_formatting(
        text: str,
    ) -> list[tuple[str, bool, bool, bool]]:
        """
        인라인 서식을 파싱하여 (텍스트, bold, italic, is_code) 튜플 리스트를 반환한다.
        """
        segments: list[tuple[str, bool, bool, bool]] = []
        # 패턴: `code`, ***bold+italic***, **bold**, *italic*
        pattern = re.compile(
            r"`([^`]+)`"              # 인라인 코드
            r"|\*\*\*(.+?)\*\*\*"    # bold+italic
            r"|\*\*(.+?)\*\*"        # bold
            r"|\*(.+?)\*"            # italic
        )

        last_end = 0
        for match in pattern.finditer(text):
            # 매치 이전 일반 텍스트
            if match.start() > last_end:
                plain = text[last_end:match.start()]
                if plain:
                    segments.append((plain, False, False, False))

            if match.group(1) is not None:
                # 인라인 코드
                segments.append((match.group(1), False, False, True))
            elif match.group(2) is not None:
                # bold+italic
                segments.append((match.group(2), True, True, False))
            elif match.group(3) is not None:
                # bold
                segments.append((match.group(3), True, False, False))
            elif match.group(4) is not None:
                # italic
                segments.append((match.group(4), False, True, False))

            last_end = match.end()

        # 나머지 텍스트
        if last_end < len(text):
            remaining = text[last_end:]
            if remaining:
                segments.append((remaining, False, False, False))

        # 빈 텍스트인 경우
        if not segments:
            segments.append((text, False, False, False))

        return segments

    @staticmethod
    def _parse_list_item(line: str) -> tuple[int, str]:
        """
        목록 항목을 파싱하여 (깊이, 정리된 텍스트)를 반환한다.
        목록이 아니면 깊이 -1을 반환한다.
        """
        # 들여쓰기 기반 깊이 계산
        stripped = line.lstrip()
        indent = len(line) - len(stripped)
        depth = indent // 2  # 2칸 들여쓰기 기준

        # 비순서 목록: -, *, +
        ul_match = re.match(r"^[-*+]\s+(.+)$", stripped)
        if ul_match:
            return depth, ul_match.group(1)

        # 순서 목록: 1., 2., ...
        ol_match = re.match(r"^\d+\.\s+(.+)$", stripped)
        if ol_match:
            return depth, ol_match.group(1)

        return -1, line

    @staticmethod
    def _set_bullet(paragraph, depth: int) -> None:
        """단락에 불릿 포인트를 설정한다."""
        pPr = paragraph._p.get_or_add_pPr()

        # 불릿 문자 설정
        bullet_char = _BULLET_CHARS[depth % len(_BULLET_CHARS)]
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", bullet_char)

        # 들여쓰기 설정
        margin_left = _INDENT_EMU * (depth + 1)
        pPr.set("marL", str(margin_left))
        pPr.set("indent", str(-_INDENT_EMU))

    @staticmethod
    def _is_code_block(line: str) -> bool:
        """코드 블록 내용인지 확인한다."""
        return line.startswith("    ") or line.startswith("\t")

    @staticmethod
    def _is_code_block_marker(line: str) -> bool:
        """코드 블록 시작/끝 마커인지 확인한다."""
        return line.strip().startswith("```")

    @staticmethod
    def _is_table_line(line: str) -> bool:
        """표 라인인지 확인한다."""
        stripped = line.strip()
        return "|" in stripped and stripped.startswith("|")

    @staticmethod
    def _extract_code_block(
        lines: list[str], start_idx: int
    ) -> tuple[str | None, int]:
        """코드 블록을 추출한다. (코드 내용, 다음 인덱스)를 반환한다."""
        if not lines[start_idx].strip().startswith("```"):
            return None, start_idx + 1

        code_lines: list[str] = []
        i = start_idx + 1
        while i < len(lines):
            if lines[i].strip() == "```":
                return "\n".join(code_lines), i + 1
            code_lines.append(lines[i])
            i += 1

        # 닫는 ``` 없이 끝난 경우
        return "\n".join(code_lines), i

    @staticmethod
    def _extract_table_data(
        lines: list[str], start_idx: int
    ) -> tuple[list[list[str]], int] | None:
        """표 데이터를 추출한다. (행 리스트, 다음 인덱스)를 반환한다."""
        rows: list[list[str]] = []
        i = start_idx

        while i < len(lines) and "|" in lines[i]:
            stripped = lines[i].strip()
            # 구분선 행 건너뛰기 (|---|---|, |:---:|---| 등)
            if re.match(r"^\|(\s*[-:]+\s*\|)+\s*$", stripped):
                i += 1
                continue
            # 셀 파싱
            cells = [
                c.strip()
                for c in stripped.split("|")
                if c.strip()
            ]
            if cells:
                rows.append(cells)
            i += 1

        if rows:
            return rows, i
        return None

    def _add_code_block(
        self,
        slide,
        code: str,
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """코드 블록을 텍스트 박스로 추가한다."""
        txBox = slide.shapes.add_textbox(left, top, width, height)

        # 배경색 설정
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = _CODE_BG_COLOR

        tf = txBox.text_frame
        tf.word_wrap = True

        code_lines = code.split("\n")
        for i, line in enumerate(code_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line
            run.font.name = self._font_config.mono_font
            run.font.size = Pt(self._font_config.code_size_pt)

    def _add_table(
        self,
        slide,
        table_data: list[list[str]],
        left: int,
        top: int,
        width: int,
        height: int,
    ) -> None:
        """표를 PPT 표 객체로 추가한다."""
        if not table_data:
            return

        rows = len(table_data)
        cols = max(len(row) for row in table_data)

        table_shape = slide.shapes.add_table(
            rows, cols, left, top, width, height
        )
        table = table_shape.table

        for row_idx, row_data in enumerate(table_data):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < cols:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = cell_text

                    # 셀 폰트 설정
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.LEFT
                        for run in paragraph.runs:
                            run.font.name = self._font_config.korean_font
                            run.font.size = Pt(self._font_config.body_size_pt - 2)

                    # 헤더 행 배경색 적용
                    if row_idx == 0:
                        cell_fill = cell.fill
                        cell_fill.solid()
                        cell_fill.fore_color.rgb = _TABLE_HEADER_BG
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = _TABLE_HEADER_FONT_COLOR
                                run.font.bold = True

    def _add_page_number(
        self, slide, template_info: TemplateInfo, page_num: int
    ) -> None:
        """슬라이드 하단에 페이지 번호를 삽입한다."""
        sw = template_info.slide_width
        sh = template_info.slide_height

        left = (sw - _PAGE_NUM_WIDTH) // 2
        top = sh - _PAGE_NUM_HEIGHT - Emu(int(_MARGIN * 0.5))

        txBox = slide.shapes.add_textbox(
            left, top, _PAGE_NUM_WIDTH, _PAGE_NUM_HEIGHT
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        # 슬라이드 번호 필드 삽입
        fld = etree.SubElement(p._p, qn("a:fld"))
        fld.set("type", "slidenum")
        fld.set("id", "{" + str(uuid.uuid4()).upper() + "}")
        rPr = etree.SubElement(fld, qn("a:rPr"))
        rPr.set("lang", "ko-KR")
        rPr.set("sz", str(int(Pt(10))))
        t_elem = etree.SubElement(fld, qn("a:t"))
        t_elem.text = str(page_num)

        # 폰트 설정
        font_elem = etree.SubElement(rPr, qn("a:latin"))
        font_elem.set("typeface", self._font_config.korean_font)

    def _format_code_paragraph(self, paragraph, line: str) -> None:
        """코드 라인을 고정폭 폰트로 포맷한다."""
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = line
        run.font.name = self._font_config.mono_font
        run.font.size = Pt(self._font_config.code_size_pt)
